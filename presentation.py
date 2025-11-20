from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Initialize Presentation
prs = Presentation()

# Define Brand Colors
BLUE_PRIMARY = RGBColor(44, 82, 130)  # #2C5282
GREEN_ACCENT = RGBColor(72, 187, 120)  # #48BB78
ORANGE_WARN = RGBColor(237, 137, 54)  # #ED8936
WHITE = RGBColor(255, 255, 255)
GRAY_BG = RGBColor(247, 250, 252)  # #F7FAFC


def add_slide(prs, layout_index, title_text, content_text=None, subtitle_text=None):
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)

    # Set Title
    title = slide.shapes.title
    title.text = title_text

    # Style Title
    title.text_frame.paragraphs[0].font.color.rgb = BLUE_PRIMARY
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(36)

    # Add Content if applicable
    if content_text:
        # If using standard content layout
        if layout_index == 1:
            body = slide.shapes.placeholders[1]
            tf = body.text_frame
            tf.text = content_text

            # Styling content
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.space_after = Pt(10)

    return slide


# --- SLIDE 1: TITLE SLIDE ---
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "AI-Powered Test Automation:\nMapping Complex Approval Workflows at Scale"
subtitle.text = "A Case Study in Human-AI Collaboration\n\nAlex Faulkner, Test Automation Lead\nDigiBlu / Experience Travel Group\nNovember 2025"

# Style Title Slide
title.text_frame.paragraphs[0].font.color.rgb = BLUE_PRIMARY
title.text_frame.paragraphs[0].font.bold = True
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)

# --- SLIDE 2: THE CHALLENGE ---
content = """The System:
• Payment system with 6 different payment types
• Up to 5 approval levels (PCM → RD → RCM → CFO → CEO)

The Complexity:
• 61 unique approval workflows
• 276 distinct rejection scenarios requiring testing

The Risk:
• Manual testing was error-prone and unmaintainable
• Business risks: Compliance violations & financial errors"""

add_slide(prs, 1, "The Challenge: Testing Complex Workflows", content)

# --- SLIDE 3: THE SCALE OF COMPLEXITY ---
content = """Why Traditional Testing Failed

The Numbers:
• 61 Unique Workflows
• 4 Raising Locations
• 6 Payment Types
• = 276 Rejection Scenarios

The Variable Logic:
• Rejection routing changes based on WHO rejected and WHEN.
• Impossible to track all permutations manually (100+ page spec)."""

add_slide(prs, 1, "The Scale of Complexity", content)

# --- SLIDE 4: THE AI SOLUTION - MAPPING ---
content = """Step 1: Ingestion
• Fed Claude AI the 100+ page Solution Spec & Excel logic.

Step 2: Pattern Recognition
• AI identified 61 unique configurations.
• Handled complex merged-cell logic in Excel.

Step 3: Output
• Dynamic Excel Workbook with 276 scenarios.
• Color-coded: Green (Happy), Red (Delete), Yellow (Reject), Blue (Recovery)."""

add_slide(prs, 1, "The AI Solution: Mapping Phase", content)

# --- SLIDE 5: THE AI SOLUTION - AUTOMATION ---
content = """From Scenarios to Scripts

• Generated ~4,700 lines of Python code (Selenium + Playwright)
• Created ~112 automated test cases supporting all 6 payment types

Advanced Features Implemented:
• PDF Validation: Downloads & verifies 11 specific fields per submission
• Email Integration: Checks approval emails via IMAP
• Evidence: Auto-captures screenshots & generates CSV reports"""

add_slide(prs, 1, "The AI Solution: Automation Phase", content)

# --- SLIDE 6: IMPLEMENTATION ---
content = """The Human + AI Collaboration Model

Humans Provided:
• Business context & domain knowledge
• Logic validation & decision making
• QA of AI outputs

AI Provided:
• Pattern recognition across 61 workflows
• Rapid code generation (~4,700 lines)
• Exhaustive scenario enumeration

Timeline: 12 Weeks from concept to 6 fully automated payment types."""

add_slide(prs, 1, "How We Implemented It", content)

# --- SLIDE 7: REAL EXAMPLE ---
content = """The Merged Cells Problem

The Challenge: 
• Excel logic used merged cells. Standard Python libs read these as "Blank".
• 50+ workflows appeared broken.

Traditional Solution:
• Manually unmerge and fix data (Est: 2-3 Weeks).

AI Solution:
• Identified issue and wrote code using 'openpyxl'.
• Logic: Detect ranges, unmerge, backfill values.
• Result: Solved in 1 day."""

add_slide(prs, 1, "Real Example: AI Solving the Impossible", content)

# --- SLIDE 8: RESULTS ---
content = """Measurable Impact

COVERAGE:
• 100% Payment Types Automated (6/6)
• 100% Rejection Scenarios Mapped (276/276)

SPEED:
• Manual Test: 2-3 hours per cycle → Automated: 5-10 minutes.
• 90-95% Time Savings per cycle.
• Development Speed: 60-70% faster with AI assistance.

QUALITY:
• Consistent Screenshot evidence & 3-tier error handling."""

add_slide(prs, 1, "The Results: By The Numbers", content)

# --- SLIDE 9: SUCCESS FACTORS ---
content = """1. Iterative Collaboration
Started small (POC) and learned from each iteration.

2. AI as Force Multiplier
Single developer + AI = Team-level output.

3. Breaking Down Complexity
Phased approach: Map Logic → Generate Scripts → Scale → Document.

4. Documentation First
AI generated docs alongside code, ensuring future-proofing and easy transfer."""

add_slide(prs, 1, "What Made This Work", content)

# --- SLIDE 10: LESSONS LEARNED ---
content = """✅ What Worked:
• Breaking complex problems into smaller pieces.
• Using AI for pattern recognition, not just coding.
• Validating AI outputs against known scenarios.

⚠️ Challenges Overcome:
• AI code was ~80% correct (needed human refinement).
• Dynamic forms required smart 'wait' logic.
• Email authentication needed specific workarounds.

💡 Key Insight:
AI accelerates execution, but Human oversight ensures correctness."""

add_slide(prs, 1, "Lessons Learned", content)

# --- SLIDE 11: BIG PICTURE ---
content = """Why This Matters & Where It Transfers

This Approach Applies To:
• Healthcare: Patient journey workflows
• Finance: Loan approval hierarchies
• Supply Chain: Procurement approvals

The Paradigm Shift:
Test automation moving from "Lagging" (months behind) to "Leading" (ready alongside features).

Core Principle:
Use AI to map complex processes that humans struggle to visualize comprehensively."""

add_slide(prs, 1, "Beyond Testing: The Big Picture", content)

# --- SLIDE 12: CLOSING ---
content = """Key Takeaways:

1. AI Democratizes Expertise (Small teams, big output).
2. Start with Understanding (Map first, code second).
3. Iterate and Validate (Human-AI loop).
4. Documentation Matters (Future-proof your solution).

What's Next:
• Expanding to self-healing tests and test data generation.

Questions?"""

add_slide(prs, 1, "Key Takeaways & Closing", content)

# Save the file
file_name = "AI_Test_Automation_Case_Study.pptx"
prs.save(file_name)
print(f"Presentation saved as {file_name}")